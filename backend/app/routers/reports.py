import io
from fastapi import APIRouter, Depends
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
import pandas as pd

from app.database import get_db
from app.models import User
from app.security import get_current_user, scope_filter
from app.services.analytics_service import query_revenue, rows_to_dicts
from app.routers.dashboard import common_filters

router = APIRouter(prefix="/api/reports", tags=["reports"])


def _get_data(db, user, filters):
    scope = scope_filter(user)
    rows = query_revenue(db, scope, **filters).all()
    return rows_to_dicts(rows)


@router.get("/csv")
def report_csv(filters: dict = Depends(common_filters), db: Session = Depends(get_db),
               user: User = Depends(get_current_user)):
    data = _get_data(db, user, filters)
    df = pd.DataFrame(data)
    buf = io.StringIO()
    df.to_csv(buf, index=False)
    buf.seek(0)
    return StreamingResponse(
        iter([buf.getvalue()]), media_type="text/csv",
        headers={"Content-Disposition": "attachment; filename=revenue_report.csv"},
    )


@router.get("/excel")
def report_excel(filters: dict = Depends(common_filters), db: Session = Depends(get_db),
                  user: User = Depends(get_current_user)):
    data = _get_data(db, user, filters)
    df = pd.DataFrame(data)
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Revenue")
        if not df.empty:
            pivot = df.pivot_table(index="customer", values="revenue", aggfunc="sum").sort_values(
                "revenue", ascending=False)
            pivot.to_excel(writer, sheet_name="By Customer")
            prod = df.pivot_table(index="product", values="revenue", aggfunc="sum").sort_values(
                "revenue", ascending=False)
            prod.to_excel(writer, sheet_name="By Product")
    buf.seek(0)
    return StreamingResponse(
        buf, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=revenue_report.xlsx"},
    )


@router.get("/pdf")
def report_pdf(filters: dict = Depends(common_filters), db: Session = Depends(get_db),
               user: User = Depends(get_current_user)):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet

    data = _get_data(db, user, filters)
    df = pd.DataFrame(data)
    total_revenue = df["revenue"].sum() if not df.empty else 0
    top = (df.groupby("customer")["revenue"].sum().sort_values(ascending=False).head(10)
           if not df.empty else pd.Series(dtype=float))

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4)
    styles = getSampleStyleSheet()
    elements = [
        Paragraph("India Post — Karnataka Circle", styles["Title"]),
        Paragraph("Revenue Intelligence Report", styles["Heading2"]),
        Spacer(1, 12),
        Paragraph(f"Total Revenue: ₹{total_revenue:,.2f} lakh", styles["Normal"]),
        Paragraph(f"Generated for: {user.username} ({user.role})", styles["Normal"]),
        Spacer(1, 16),
        Paragraph("Top 10 Customers by Revenue", styles["Heading3"]),
    ]
    table_data = [["#", "Customer", "Revenue (Lakh)"]] + [
        [str(i + 1), name, f"{val:,.2f}"] for i, (name, val) in enumerate(top.items())
    ]
    t = Table(table_data, hAlign="LEFT")
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#B01E28")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
    ]))
    elements.append(t)
    doc.build(elements)
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
                              headers={"Content-Disposition": "attachment; filename=revenue_report.pdf"})


@router.get("/ppt")
def report_ppt(filters: dict = Depends(common_filters), db: Session = Depends(get_db),
               user: User = Depends(get_current_user)):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    data = _get_data(db, user, filters)
    df = pd.DataFrame(data)
    total_revenue = df["revenue"].sum() if not df.empty else 0
    top = (df.groupby("customer")["revenue"].sum().sort_values(ascending=False).head(8)
           if not df.empty else pd.Series(dtype=float))

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "India Post — Karnataka Circle"
    slide.placeholders[1].text = "Revenue Intelligence Summary"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body = slide2.placeholders[1].text_frame
    body.text = f"Total Revenue: ₹{total_revenue:,.2f} lakh"
    for name, val in top.items():
        p = body.add_paragraph()
        p.text = f"{name}: ₹{val:,.2f} lakh"
        p.level = 1

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=revenue_report.pptx"},
    )
